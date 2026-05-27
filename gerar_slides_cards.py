"""
gerar_slides_cards.py — PPT: 4 cartoes v4.0 + 4 slides tecnicos do CRITICO
Uso: python gerar_slides_cards.py
Saida: apresentacao_fb14_cards.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── PALETA ───────────────────────────────────────────────────────────────────

SLIDE_BG = RGBColor(0x14, 0x14, 0x1E)
SEC_A    = RGBColor(0x2A, 0x2A, 0x38)
SEC_B    = RGBColor(0x1E, 0x1E, 0x2C)
BAR_BG   = RGBColor(0x44, 0x44, 0x55)
WHT      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x88, 0x88, 0x99)
BLK      = RGBColor(0x00, 0x00, 0x00)
ACCENT   = RGBColor(0x55, 0xAA, 0xFF)

HEADER_COLORS = {
    "RISCO":       (RGBColor(0xB8, 0x80, 0x00), BLK),
    "AVISO":       (RGBColor(0x00, 0x5C, 0xA8), WHT),
    "CONFIRMADO":  (RGBColor(0x8B, 0x1A, 0x1A), WHT),
    "FIM_DE_VIDA": (RGBColor(0x16, 0x16, 0x16), WHT),
}
BAR_FILL = {
    "AVISO":       RGBColor(0x00, 0x5C, 0xA8),
    "CONFIRMADO":  RGBColor(0xAA, 0x20, 0x20),
    "FIM_DE_VIDA": RGBColor(0xDD, 0x20, 0x20),
}
LEFT_LABEL = {
    "RISCO":       RGBColor(0xD4, 0x9A, 0x00),
    "AVISO":       RGBColor(0x44, 0x99, 0xFF),
    "CONFIRMADO":  RGBColor(0xFF, 0x55, 0x55),
    "FIM_DE_VIDA": RGBColor(0xFF, 0x66, 0x00),
}

# ─── LAYOUT ───────────────────────────────────────────────────────────────────

SW, SH = Inches(13.33), Inches(7.5)
CX, CW = Inches(4.55), Inches(4.4)
CT      = Inches(0.55)
PX, PY  = Inches(0.14), Inches(0.10)
GAP     = Inches(0.05)
ROW_H   = Inches(0.32)
LBL_H   = Inches(0.22)

# ─── PRIMITIVAS ───────────────────────────────────────────────────────────────

def _bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = SLIDE_BG

def _r(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _t(slide, x, y, w, h, text, pt, bold=False, col=WHT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(pt)
    run.font.bold = bold
    run.font.color.rgb = col
    return tb

def _tlines(slide, x, y, w, h, lines, pt, bold=False, col=WHT):
    """Multiline textbox: lines is list of strings."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(pt)
        run.font.bold = bold
        run.font.color.rgb = col
    return tb

# ─── SECOES DO CARD ───────────────────────────────────────────────────────────

def _header(slide, tipo, titulo, subtitulo, maquina, data_str):
    H = Inches(1.10)
    hbg, htxt = HEADER_COLORS[tipo]
    _r(slide, CX, CT, CW, H, hbg)
    _t(slide, CX+PX,       CT+Inches(0.07), CW*0.56, Inches(0.45), titulo,   12, bold=True, col=htxt)
    _t(slide, CX+PX,       CT+Inches(0.52), CW*0.56, Inches(0.26), subtitulo, 8, col=htxt)
    _t(slide, CX+CW*0.55,  CT+Inches(0.07), CW*0.43, Inches(0.42), maquina,  14, bold=True, col=htxt, align=PP_ALIGN.RIGHT)
    _t(slide, CX+CW*0.55,  CT+Inches(0.53), CW*0.43, Inches(0.22), data_str,  8, col=htxt, align=PP_ALIGN.RIGHT)
    return CT + H + GAP

def _vida(slide, tipo, y, dias_op, frac, dias_rest, alem_eta=False):
    H = Inches(1.05)
    _r(slide, CX, y, CW, H, SEC_A)
    _t(slide, CX+PX, y+PY, CW-2*PX, LBL_H, "VIDA DO MAINTACKER", 8, bold=True, col=GRAY)
    bx, by = CX+PX, y+Inches(0.33)
    bw, bh = CW-2*PX-Inches(0.92), Inches(0.20)
    _r(slide, bx, by, bw, bh, BAR_BG)
    fw = max(int(bw * min(frac, 1.0)), int(Inches(0.05)))
    _r(slide, bx, by, fw, bh, BAR_FILL.get(tipo, BAR_FILL["CONFIRMADO"]))
    _t(slide, bx+bw+Inches(0.06), by-Inches(0.01), Inches(0.85), Inches(0.24),
       f"{round(frac*100)}% consumida", 8, bold=True)
    _t(slide, CX+PX, y+Inches(0.62), CW*0.55, Inches(0.22),
       f"{dias_op} dias em operacao", 8, col=GRAY)
    lbl = f"{dias_rest} dias alem do ETA" if alem_eta else f"~{dias_rest} dias para troca"
    _t(slide, CX+CW*0.55, y+Inches(0.62), CW*0.43, Inches(0.22), lbl, 8, col=GRAY, align=PP_ALIGN.RIGHT)
    return y + H + GAP

def _section(slide, y, label, facts, bg=SEC_B):
    n = len(facts)
    H = PY + LBL_H + Inches(0.06) + n*ROW_H + Inches(0.10)
    _r(slide, CX, y, CW, H, bg)
    _t(slide, CX+PX, y+PY, CW-2*PX, LBL_H, label, 8, bold=True, col=GRAY)
    ry = y + PY + LBL_H + Inches(0.06)
    for title, value in facts:
        _t(slide, CX+PX,          ry, Inches(2.05), ROW_H, title, 9, bold=True)
        _t(slide, CX+Inches(2.25), ry, CW-Inches(2.25)-PX, ROW_H, value, 9)
        ry += ROW_H
    return y + H + GAP

def _acao(slide, y, text):
    H = Inches(1.05)
    _r(slide, CX, y, CW, H, SEC_A)
    _t(slide, CX+PX, y+PY, CW-2*PX, LBL_H, "ACAO RECOMENDADA", 8, bold=True, col=GRAY)
    _t(slide, CX+PX, y+PY+LBL_H+Inches(0.05), CW-2*PX, Inches(0.68), text, 9, bold=True)

def _left_panel(slide, card):
    tipo = card["tipo"]
    lx, lw = Inches(0.4), CX - Inches(0.7)
    _t(slide, lx, Inches(0.85), lw, Inches(0.3),  f"Dia {card['dia']} do ciclo", 10, col=GRAY)
    _t(slide, lx, Inches(1.25), lw, Inches(0.75), tipo.replace("_", " "), 26, bold=True, col=LEFT_LABEL[tipo])
    _t(slide, lx, Inches(2.10), lw, Inches(1.10), card["descricao"], 10, col=GRAY)
    _t(slide, lx, Inches(3.55), lw, Inches(0.24), "QUANDO DISPARA", 8, bold=True, col=GRAY)
    _t(slide, lx, Inches(3.85), lw, Inches(0.85), card["quando"], 10, col=WHT)
    hbg, _ = HEADER_COLORS[tipo]
    _r(slide, lx, Inches(5.15), Inches(0.22), Inches(0.22), hbg)
    _t(slide, lx+Inches(0.32), Inches(5.13), lw-Inches(0.35), Inches(0.28), card["cor_desc"], 9, col=GRAY)

# ─── DADOS DOS CARTOES ────────────────────────────────────────────────────────

CARDS = [
    {
        "tipo": "RISCO", "dia": 8,
        "titulo": "RISCO — Leitura Anomala",
        "subtitulo": "Forca de Selagem — Rolo Maintacker",
        "maquina": "FB14-TESTE", "data_str": "25/04/2026 06:00",
        "descricao": "Leitura anomala isolada. Forca abaixo de 800 gf sem padrao de degradacao ou envelhecimento.",
        "quando": "Forca < 800 gf, mediana do sinal > 950 gf, sem evidencia de degradacao.",
        "cor_desc": "Amarelo — alerta pontual, sem urgencia imediata",
        "evento_facts": [
            ("Forca minima registrada", "732 gf"),
            ("Forca media (3 dias)",    "1104 gf"),
            ("Data do evento",          "25/04/2026"),
            ("< 800 gf no ciclo",       "1x"),
        ],
        "dias_op_str": "8 dias",
        "acao": "Ir ao local e verificar os motivos da forca de selagem abaixo de 800 gf. Rolo novo — sem degradacao associada.",
    },
    {
        "tipo": "AVISO", "dia": 20,
        "titulo": "AVISO — Degradacao Detectada",
        "subtitulo": "Forca de Selagem — Rolo Maintacker",
        "maquina": "FB14-TESTE", "data_str": "07/05/2026 06:00",
        "descricao": "Degradacao precoce detectada. Tendencia de queda confirmada com probabilidade de risco moderada.",
        "quando": "p_risk >= 35% OU signal_score > 15%",
        "cor_desc": "Azul — degradacao iniciada, monitoramento reforcado",
        "vida_frac": 20/54, "dias_op": 20, "dias_rest": 34, "alem_eta": False,
        "indicadores": [
            ("Tendencia de Forca",      "Declinio leve (-5 gf/dia)"),
            ("Forca Projetada (48h)",   "900 gf"),
            ("Forca Minima (3 dias)",   "920 gf"),
            ("Forca Minima do Ciclo",   "732 gf"),
            ("< 800 gf no ciclo atual", "1x"),
            ("Media da Semana Atual",   "940 gf"),
            ("Media da Semana Passada", "1050 gf"),
        ],
        "acao": "Aumentar frequencia de monitoramento do sinal de forca. Registrar observacoes no proximo turno.",
    },
    {
        "tipo": "CONFIRMADO", "dia": 30,
        "titulo": "ALERTA CRITICO",
        "subtitulo": "Forca de Selagem — Rolo Maintacker",
        "maquina": "FB14-TESTE", "data_str": "17/05/2026 06:00",
        "descricao": "Degradacao severa confirmada. Forca critica registrada com multiplos eventos criticos no ciclo.",
        "quando": "Forca < 800 gf + p_risk >= 40% + 2 ou mais eventos criticos no ciclo",
        "cor_desc": "Vermelho — intervencao preventiva necessaria",
        "vida_frac": 30/54, "dias_op": 30, "dias_rest": 24, "alem_eta": False,
        "indicadores": [
            ("Tendencia de Forca",      "Declinio acentuado (-11 gf/dia)"),
            ("Forca Projetada (48h)",   "755 gf"),
            ("Forca Minima (3 dias)",   "778 gf"),
            ("Forca Minima do Ciclo",   "732 gf"),
            ("< 800 gf no ciclo atual", "4x"),
            ("Media da Semana Atual",   "820 gf"),
            ("Media da Semana Passada", "960 gf"),
        ],
        "acao": "Analise aprofundada e planejamento de troca do rolo maintacker. Forca critica confirmada — agendar intervencao preventiva.",
    },
    {
        "tipo": "FIM_DE_VIDA", "dia": 49,
        "titulo": "FIM DE VIDA — Troca Imediata",
        "subtitulo": "Forca de Selagem — Rolo Maintacker",
        "maquina": "FB14-TESTE", "data_str": "05/06/2026 06:00",
        "descricao": "Vida util projetada atingida. Rolo dentro da janela critica de 5 dias antes do ETA de 54 dias.",
        "quando": "Idade >= ETA - 5 dias (janela de fim de vida Weibull)",
        "cor_desc": "Escuro — troca imediata, sem margem operacional",
        "vida_frac": 49/54, "dias_op": 49, "dias_rest": 5, "alem_eta": False,
        "indicadores": [
            ("Tendencia de Forca",      "Declinio acentuado (-14 gf/dia)"),
            ("Forca Projetada (48h)",   "715 gf"),
            ("Forca Minima (3 dias)",   "741 gf"),
            ("Forca Minima do Ciclo",   "698 gf"),
            ("< 800 gf no ciclo atual", "7x"),
            ("Media da Semana Atual",   "780 gf"),
            ("Media da Semana Passada", "850 gf"),
        ],
        "acao": "Troca imediata do rolo maintacker — vida util projetada atingida. ETA: 54 dias. Rolo com 49 dias em operacao.",
    },
]

# ─── CONSTRUTORES DE SLIDE ────────────────────────────────────────────────────

def build_risco(prs, card):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _left_panel(slide, card)
    _header(slide, "RISCO", card["titulo"], card["subtitulo"], card["maquina"], card["data_str"])
    y = CT + Inches(1.10) + GAP
    y = _section(slide, y, "EVENTO",         card["evento_facts"],                              bg=SEC_A)
    y = _section(slide, y, "STATUS DO ROLO", [("Dias em operacao", card["dias_op_str"])],        bg=SEC_B)
    _acao(slide, y, card["acao"])


def build_critico(prs, card):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _left_panel(slide, card)
    y = _header(slide, card["tipo"], card["titulo"], card["subtitulo"], card["maquina"], card["data_str"])
    y = _vida(slide, card["tipo"], y, card["dias_op"], card["vida_frac"], card["dias_rest"], card["alem_eta"])
    y = _section(slide, y, "INDICADORES", card["indicadores"])
    _acao(slide, y, card["acao"])


# ─── SLIDES TECNICOS ─────────────────────────────────────────────────────────

def _tech_header(slide, title):
    _t(slide, Inches(0.5), Inches(0.22), SW-Inches(1.0), Inches(0.5), title, 20, bold=True, col=ACCENT)
    _r(slide, Inches(0.5), Inches(0.76), SW-Inches(1.0), Inches(0.025), RGBColor(0x33, 0x55, 0x88))


def tech_anatomia(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _tech_header(slide, "ANATOMIA DO CARTAO CRITICO — VISAO GERAL")

    cols = [
        ("1", "CABECALHO",          HEADER_COLORS["CONFIRMADO"][0],
         "Identifica nivel de alerta (AVISO / CONFIRMADO / FIM DE VIDA), maquina monitorada e timestamp exato do disparo."),
        ("2", "VIDA DO MAINTACKER", SEC_A,
         "Barra visual baseada no modelo Weibull (eta=54d). Mostra % de vida consumida e dias restantes estimados."),
        ("3", "INDICADORES",        SEC_B,
         "7 metricas em tempo real: tendencia de forca, projecoes 48h, minimos, contagem de eventos criticos e medias semanais."),
        ("4", "ACAO RECOMENDADA",   SEC_A,
         "Protocolo especifico para o nivel disparado: do reforco de monitoramento ate a troca imediata."),
    ]

    col_w = (SW - Inches(1.0)) / 4
    for i, (num, title, color, desc) in enumerate(cols):
        cx = Inches(0.5) + i * col_w
        _r(slide, cx+Inches(0.05), Inches(1.0), col_w-Inches(0.1), Inches(0.50), color)
        _t(slide, cx+Inches(0.12), Inches(1.07), col_w-Inches(0.2), Inches(0.38),
           f"{num}. {title}", 10, bold=True)
        _r(slide, cx+Inches(0.05), Inches(1.52), col_w-Inches(0.1), Inches(5.2), SEC_B)
        _t(slide, cx+Inches(0.12), Inches(1.6),  col_w-Inches(0.2), Inches(2.0), desc, 10, col=GRAY)

    _t(slide, Inches(0.5), SH-Inches(0.55), SW-Inches(1.0), Inches(0.38),
       "Cada secao e gerada dinamicamente pelo trigger_engine.py com dados em tempo real do Seeq Data Lab.",
       9, col=GRAY)


def tech_vida(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _tech_header(slide, "BARRA DE VIDA — MODELO WEIBULL")

    stages = [
        ("Ciclo novo  (Dia  8)",  8/54,  HEADER_COLORS["RISCO"][0],  "15%"),
        ("AVISO       (Dia 20)", 20/54,  BAR_FILL["AVISO"],           "37%"),
        ("CONFIRMADO  (Dia 30)", 30/54,  BAR_FILL["CONFIRMADO"],      "56%"),
        ("FIM DE VIDA (Dia 49)", 49/54,  BAR_FILL["FIM_DE_VIDA"],     "91%"),
    ]

    bx, bw, bh = Inches(0.5), Inches(8.2), Inches(0.26)
    for i, (label, frac, col, pct) in enumerate(stages):
        yr = Inches(1.1) + i * Inches(1.35)
        _t(slide, bx, yr, Inches(2.6), Inches(0.28), label, 10, col=WHT)
        by = yr + Inches(0.30)
        _r(slide, bx, by, bw, bh, BAR_BG)
        _r(slide, bx, by, int(bw * frac), bh, col)
        _t(slide, bx+bw+Inches(0.10), by-Inches(0.01), Inches(0.82), Inches(0.28),
           pct, 10, bold=True, col=col)

    # Weibull box on right
    rx, ry = Inches(9.1), Inches(0.95)
    rw = SW - rx - Inches(0.3)
    _r(slide, rx, ry, rw, Inches(5.9), SEC_B)
    _t(slide, rx+PX, ry+PY, rw-2*PX, Inches(0.26), "MODELO WEIBULL", 9, bold=True, col=ACCENT)
    _tlines(slide, rx+PX, ry+Inches(0.44), rw-2*PX, Inches(5.2), [
        "Beta (beta): 1.181",
        "Eta (eta):   1297 h = 54 dias",
        "",
        "consumida = idade / eta_ajustado",
        "",
        "eta_ajustado = eta - desvio * 0.8",
        "  desvio = eta - media_ano_vigente",
        "",
        "p_risk = F(t | beta, eta)",
        "  acumula com a idade",
        "  + evidencia dos sinais de forca",
    ], 9, col=WHT)


def tech_indicadores(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _tech_header(slide, "OS 7 INDICADORES DO CARTAO CRITICO")

    headers  = ["Campo",                 "Dia 30",   "Como e calculado",                          "Limiar"]
    col_ws   = [Inches(2.55), Inches(1.30), Inches(5.60), Inches(3.50)]
    col_xs   = []
    cx0 = Inches(0.28)
    for j, w in enumerate(col_ws):
        col_xs.append(cx0 + sum(col_ws[:j]))

    hy = Inches(0.90)
    for hdr, cw, cx in zip(headers, col_ws, col_xs):
        _r(slide, cx, hy, cw-Inches(0.03), Inches(0.30), RGBColor(0x22, 0x44, 0x77))
        _t(slide, cx+Inches(0.05), hy+Inches(0.03), cw, Inches(0.25), hdr, 9, bold=True)

    rows = [
        ("Tendencia de Forca",       "Decl. acentuado\n(-11 gf/dia)", "Regressao linear sobre medias diarias dos ultimos 7 dias",              "> -10 gf/dia = critico"),
        ("Forca Projetada (48h)",    "755 gf",  "Extrapolacao linear: ultima media + slope * 2 dias",                          "< 800 gf = alerta"),
        ("Forca Minima (3 dias)",    "778 gf",  "Menor valor absoluto nos ultimos 3 dias de operacao",                        "< 800 gf = risco"),
        ("Forca Minima do Ciclo",    "732 gf",  "Menor valor absoluto desde a ultima troca do rolo",                          "< 800 gf = historico"),
        ("< 800 gf no ciclo atual",  "4x",       "Contagem de leituras abaixo de 800 gf desde a troca",                      ">= 3x = critico"),
        ("Media da Semana Atual",    "820 gf",  "Media das leituras dos ultimos 7 dias",                                      "< 950 gf = degradacao"),
        ("Media da Semana Passada",  "960 gf",  "Media dos 7 dias anteriores — base de comparacao para calcular a queda",   "Comparativo de queda"),
    ]

    rh = Inches(0.755)
    for i, (campo, valor, calc, limiar) in enumerate(rows):
        yr = Inches(1.25) + i * rh
        bg_row = SEC_A if i % 2 == 0 else SEC_B
        for cx, cw in zip(col_xs, col_ws):
            _r(slide, cx, yr, cw-Inches(0.03), rh-Inches(0.02), bg_row)
        vals = [campo, valor, calc, limiar]
        for j, (val, cx, cw) in enumerate(zip(vals, col_xs, col_ws)):
            _t(slide, cx+Inches(0.06), yr+Inches(0.05), cw-Inches(0.1), rh-Inches(0.06),
               val, 8.5, bold=(j == 0))


def tech_protocolo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _tech_header(slide, "PROTOCOLO DE RESPOSTA POR NIVEL DE ALERTA")

    niveis = [
        {
            "tipo": "AVISO", "hcol": HEADER_COLORS["AVISO"][0],
            "urgencia": "MODERADA", "prazo": "24 – 48 h",
            "acoes": [
                "Aumentar frequencia de monitoramento",
                "Verificar condicoes operacionais do turno",
                "Registrar observacoes no sistema",
                "Preparar plano de intervencao preventiva",
            ],
            "condicao": "p_risk >= 35%  OU  signal_score > 15%",
        },
        {
            "tipo": "CONFIRMADO", "hcol": HEADER_COLORS["CONFIRMADO"][0],
            "urgencia": "ALTA", "prazo": "Proximo turno",
            "acoes": [
                "Analise aprofundada do sinal de forca",
                "Agendar troca preventiva do rolo",
                "Confirmar disponibilidade de pecas",
                "Notificar lideranca de manutencao",
            ],
            "condicao": "Forca < 800 gf  +  p_risk >= 40%  +  >= 2 eventos no ciclo",
        },
        {
            "tipo": "FIM DE VIDA", "hcol": HEADER_COLORS["FIM_DE_VIDA"][0],
            "urgencia": "CRITICA", "prazo": "Imediato",
            "acoes": [
                "Troca imediata do rolo maintacker",
                "Parar producao se necessario",
                "Executar checklist de troca",
                "Registrar data de troca no sistema",
            ],
            "condicao": "Idade >= ETA - 5 dias  (janela Weibull fim de vida)",
        },
    ]

    col_w = (SW - Inches(1.0)) / 3
    for i, n in enumerate(niveis):
        cx = Inches(0.5) + i * col_w
        cy = Inches(1.05)
        _r(slide, cx+Inches(0.05), cy, col_w-Inches(0.10), Inches(0.50), n["hcol"])
        _t(slide, cx+Inches(0.12), cy+Inches(0.07), col_w-Inches(0.2), Inches(0.38),
           n["tipo"], 13, bold=True)

        uy = cy + Inches(0.58)
        _r(slide, cx+Inches(0.05), uy, col_w-Inches(0.10), Inches(0.60), SEC_B)
        _t(slide, cx+Inches(0.12), uy+Inches(0.04), col_w*0.45, Inches(0.20), "URGENCIA", 7, bold=True, col=GRAY)
        _t(slide, cx+Inches(0.12), uy+Inches(0.26), col_w*0.45, Inches(0.26), n["urgencia"], 10, bold=True, col=n["hcol"])
        _t(slide, cx+col_w*0.5,  uy+Inches(0.04), col_w*0.47, Inches(0.20), "PRAZO",    7, bold=True, col=GRAY)
        _t(slide, cx+col_w*0.5,  uy+Inches(0.26), col_w*0.47, Inches(0.26), n["prazo"], 10, bold=True, col=WHT)

        ay = uy + Inches(0.68)
        _r(slide, cx+Inches(0.05), ay, col_w-Inches(0.10), Inches(2.55), SEC_A)
        _t(slide, cx+Inches(0.12), ay+Inches(0.06), col_w-Inches(0.2), Inches(0.22), "ACOES", 8, bold=True, col=GRAY)
        for j, acao in enumerate(n["acoes"]):
            _t(slide, cx+Inches(0.12), ay+Inches(0.32)+j*Inches(0.54),
               col_w-Inches(0.22), Inches(0.48), f"• {acao}", 9)

        ty = ay + Inches(2.63)
        _r(slide, cx+Inches(0.05), ty, col_w-Inches(0.10), Inches(0.55), SEC_B)
        _t(slide, cx+Inches(0.12), ty+Inches(0.04), col_w-Inches(0.2), Inches(0.20),
           "CONDICAO DE DISPARO", 7, bold=True, col=GRAY)
        _t(slide, cx+Inches(0.12), ty+Inches(0.26), col_w-Inches(0.2), Inches(0.24),
           n["condicao"], 8, col=n["hcol"])


# ─── MAIN ─────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    build_risco(prs,   CARDS[0])
    build_critico(prs, CARDS[1])
    build_critico(prs, CARDS[2])
    build_critico(prs, CARDS[3])

    tech_anatomia(prs)
    tech_vida(prs)
    tech_indicadores(prs)
    tech_protocolo(prs)

    out = "apresentacao_fb14_cards.pptx"
    prs.save(out)
    print(f"Salvo: {out}  ({len(prs.slides)} slides)")
